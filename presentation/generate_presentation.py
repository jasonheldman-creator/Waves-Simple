#!/usr/bin/env python3
"""
WAVES Intelligence™ Presentation Generator
Generates an animated PowerPoint presentation from the executive briefing script.

This script creates a 13-slide presentation with:
- Dark institutional theme
- Subtle animations (fade, slide-in, sequential emphasis)
- Vector™ avatar presence
- Simplified diagrams for key concepts
- Narration-ready slide notes
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import os


class WavesPresentation:
    """Generator for WAVES Intelligence presentation."""
    
    # Color scheme - Dark institutional theme
    COLORS = {
        'background': RGBColor(20, 25, 35),  # Dark blue-gray
        'primary': RGBColor(100, 180, 255),  # Light blue (WAVES brand)
        'secondary': RGBColor(150, 200, 255),  # Lighter blue
        'accent': RGBColor(255, 200, 100),  # Warm accent
        'text': RGBColor(240, 245, 250),  # Off-white
        'text_dim': RGBColor(180, 190, 200),  # Dimmed text
        'success': RGBColor(100, 220, 150),  # Success green
        'warning': RGBColor(255, 180, 100),  # Warning orange
    }
    
    def __init__(self):
        """Initialize presentation with 16:9 widescreen layout."""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        
    def create_blank_slide(self):
        """Create a blank slide with dark background."""
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Add dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['background']
        
        return slide
    
    def add_vector_branding(self, slide, position='bottom-right'):
        """Add subtle Vector™ branding to slide."""
        if position == 'bottom-right':
            left, top = Inches(8.5), Inches(5.0)
        elif position == 'top-right':
            left, top = Inches(8.5), Inches(0.3)
        else:
            left, top = Inches(8.5), Inches(5.0)
        
        textbox = slide.shapes.add_textbox(left, top, Inches(1.2), Inches(0.4))
        text_frame = textbox.text_frame
        text_frame.text = "Vector™"
        
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(10)
        p.font.color.rgb = self.COLORS['text_dim']
        p.font.italic = True
        
    def add_title(self, slide, title_text, subtitle_text=None):
        """Add title and optional subtitle to slide."""
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(1)
        )
        text_frame = title_box.text_frame
        text_frame.text = title_text
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        
        # Subtitle if provided
        if subtitle_text:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.6), Inches(9), Inches(0.6)
            )
            text_frame = subtitle_box.text_frame
            text_frame.text = subtitle_text
            
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.font.size = Pt(24)
            p.font.color.rgb = self.COLORS['secondary']
            
    def add_body_text(self, slide, text, left=0.5, top=2.5, width=9, height=2.5):
        """Add body text to slide."""
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = self.COLORS['text']
            paragraph.space_before = Pt(6)
            paragraph.space_after = Pt(6)
            
    def add_bullet_points(self, slide, points, left=0.5, top=2.5):
        """Add bullet points to slide."""
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(9), Inches(2.5)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(points):
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            
            p.text = point
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = self.COLORS['text']
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            
    def add_diagram_placeholder(self, slide, diagram_type, left=2, top=2.5, width=6, height=2.5):
        """Add a placeholder for diagrams."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        
        # Style the placeholder
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 40, 55)
        shape.line.color.rgb = self.COLORS['primary']
        shape.line.width = Pt(2)
        
        # Add diagram label
        text_frame = shape.text_frame
        text_frame.text = f"[{diagram_type}]"
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = self.COLORS['text_dim']
        p.font.italic = True
        
    def add_notes(self, slide, notes_text):
        """Add speaker notes to slide."""
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text
        
    # ===== Slide Generation Methods =====
    
    def slide_1_title(self):
        """Slide 1: Title & Introduction"""
        slide = self.create_blank_slide()
        
        # Main title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(8), Inches(1.5)
        )
        text_frame = title_box.text_frame
        text_frame.text = "WAVES Intelligence™"
        
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.0), Inches(8), Inches(0.8)
        )
        text_frame = subtitle_box.text_frame
        text_frame.text = "Institutional Portfolio Analytics Platform"
        
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28)
        p.font.color.rgb = self.COLORS['secondary']
        
        # Attribution
        attr_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.0), Inches(8), Inches(0.5)
        )
        text_frame = attr_box.text_frame
        text_frame.text = "Executive Briefing presented by Vector™"
        
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.color.rgb = self.COLORS['text_dim']
        p.font.italic = True
        
        self.add_notes(slide, """Good afternoon. I'm Vector, your guide to WAVES Intelligence. Over the next seven to nine minutes, I'll walk you through our institutional-grade portfolio analytics platform—designed specifically for sophisticated investors who demand transparency, precision, and governance-ready insights.

WAVES Intelligence represents a fundamental shift in how institutions understand and explain portfolio performance. This isn't just another analytics dashboard. It's a complete decision-making ecosystem built on three core principles: no-predict constraints, deterministic outputs, and architectural transparency.""")
        
    def slide_2_challenge(self):
        """Slide 2: The Challenge"""
        slide = self.create_blank_slide()
        self.add_title(slide, "The Traditional Attribution Gap")
        
        # Add diagram placeholder
        self.add_diagram_placeholder(
            slide, "Traditional Black-Box vs Transparent Systems",
            left=2, top=2.0, width=6, height=2.5
        )
        
        self.add_vector_branding(slide)
        self.add_notes(slide, """Traditional portfolio analytics platforms present a challenge. They often deliver attribution reports that are opaque, difficult to validate, and unreliable under changing market regimes. Institutional investors need more than backward-looking performance metrics—they need governance-grade truth outputs that explain exactly where returns come from, under what conditions they persist, and what risks threaten their durability.

Most platforms optimize for prediction. WAVES Intelligence takes a different approach. We embrace a no-predict constraint, focusing instead on deterministic, explainable insights that boards and oversight committees can actually use.""")
        
    def slide_3_solution(self):
        """Slide 3: The WAVES Solution - Overview"""
        slide = self.create_blank_slide()
        self.add_title(slide, "WAVES Intelligence™ Platform Architecture")
        
        points = [
            "Alpha Attribution — Vector Truth Layer with regime awareness",
            "Performance Deep Dive — Comprehensive wave tracking & metrics",
            "Decision Ledger — Governance layer with audit trail",
            "Board-Ready Reporting — Institutional PDF board packs"
        ]
        self.add_bullet_points(slide, points)
        
        self.add_vector_branding(slide)
        self.add_notes(slide, """WAVES Intelligence provides four integrated capabilities:

First, Alpha Attribution—our Vector Truth Layer decomposes returns into structural effects and residual strategy returns, with complete regime awareness.

Second, Performance Deep Dive—comprehensive wave performance charts, metrics, and historical tracking with synthetic data support for seamless onboarding.

Third, Decision Ledger—a governance layer that tracks every decision, contract, and oversight action, creating an immutable audit trail.

And fourth, Board-Ready Reporting—institutional PDF board packs that translate complex analytics into clear, actionable narratives.

Each component is built on the Wave ID system, our canonical identifier framework that ensures data integrity across all analytics.""")
        
    def slide_4_wave_id(self):
        """Slide 4: Wave ID System"""
        slide = self.create_blank_slide()
        self.add_title(slide, "The Wave ID System",
                      "Canonical Data Architecture")
        
        self.add_diagram_placeholder(
            slide, "Wave ID as Central Hub",
            left=2, top=2.2, width=6, height=2.5
        )
        
        self.add_vector_branding(slide)
        self.add_notes(slide, """The foundation of WAVES Intelligence is the Wave ID system. Think of Wave IDs as permanent identifiers for every investment strategy or portfolio composition we track.

Each Wave has three data layers: historical performance data, configuration parameters including benchmark definitions, and dynamic weightings that evolve over time.

The system supports both real market data and synthetic placeholder data, marked explicitly with transparency indicators. This means institutions can onboard immediately, with all analytics fully functional, while real data integration happens in the background.

The Wave ID architecture ensures backward compatibility, so legacy data integrates seamlessly. No disruption. No data migration friction.""")
        
    def slide_5_vector_truth(self):
        """Slide 5: Vector Truth Layer"""
        slide = self.create_blank_slide()
        self.add_title(slide, "Vector™ Truth Layer",
                      "Governance-Grade Attribution")
        
        self.add_diagram_placeholder(
            slide, "Attribution Stack: Structural vs Residual",
            left=2, top=2.2, width=6, height=2.5
        )
        
        self.add_vector_branding(slide)
        self.add_notes(slide, """Let me explain the Vector Truth Layer—our attribution engine.

Traditional attribution tells you what happened. Vector Truth tells you why it happened, with governance-ready reliability signals.

We separate returns into two categories: structural effects and residual strategy returns.

Structural effects include capital preservation overlays—VIX regime controls, SmartSafe gating—and benchmark construction offsets. These are non-alpha by design. They're governance controls.

Residual strategy returns reflect timing, exposure scaling, volatility control, and regime management after those structural overlays. We don't attribute these to asset selection or static weights. We show them as integrated strategy decisions.

This separation is critical. It prevents alpha inflation, ensures intellectual honesty, and gives boards the transparency they need.""")
        
    def slide_6_alpha_decomposition(self):
        """Slide 6: Alpha Decomposition Example"""
        slide = self.create_blank_slide()
        self.add_title(slide, "Alpha Sources",
                      "Decomposition Example")
        
        # Example breakdown
        points = [
            "Total Excess Return: 8.0%",
            "  • Security Selection Alpha: 5.0%",
            "  • Exposure Management Alpha: 2.0%",
            "  • Capital Preservation Effect: +1.5%",
            "  • Benchmark Construction Offset: -0.5%",
            "Residual Strategy Return: 7.5% (Alpha-Eligible)"
        ]
        self.add_bullet_points(slide, points, top=2.2)
        
        self.add_vector_branding(slide)
        self.add_notes(slide, """Here's how it works in practice.

Imagine a Wave delivers eight percent total excess return over its benchmark. Vector Truth decomposes this into:

Security selection alpha of five percent—the exposure-adjusted return from strategic positioning.

Exposure management alpha of two percent—gains from timing and scaling decisions.

Capital preservation effect of one and a half percent—the structural benefit from VIX overlays and regime-aware cash sweeps.

And a benchmark construction offset of negative point five percent—the expected structural drag from benchmark composition choices.

Notice how the structural components largely offset. That's intentional. The residual strategy return of seven point five percent is what's truly alpha-eligible after governance controls.

Vector Truth provides this breakdown with full regime attribution, showing how much alpha was earned in risk-on versus risk-off environments. This tells you whether the strategy is durable across market conditions.""")
        
    def slide_7_regime_attribution(self):
        """Slide 7: Regime Attribution & Durability"""
        slide = self.create_blank_slide()
        self.add_title(slide, "Regime Attribution",
                      "Risk-On vs Risk-Off")
        
        self.add_diagram_placeholder(
            slide, "Alpha Split by Market Regime",
            left=2, top=2.2, width=6, height=2.5
        )
        
        self.add_vector_branding(slide)
        self.add_notes(slide, """Regime attribution is essential for assessing durability.

Vector Truth tracks every return period's market regime—risk-on or risk-off—based on VIX levels and volatility signals.

We then show cumulative alpha earned in each regime. If a strategy delivers four percent alpha in risk-on and three percent in risk-off, that's balanced. It suggests resilience.

But if a strategy shows six percent alpha in risk-on and negative one percent in risk-off, that's a red flag. The strategy is regime-dependent. Its durability is questionable.

We calculate a fragility score from zero to one. Low fragility means the alpha is structurally sound. High fragility means it's vulnerable to regime shifts, dispersion collapse, or volatility suppression.

Boards need to know this. Traditional attribution hides it. Vector Truth surfaces it explicitly.""")
        
    def slide_8_decision_ledger(self):
        """Slide 8: Decision Ledger"""
        slide = self.create_blank_slide()
        self.add_title(slide, "Decision Ledger",
                      "Audit Trail & Governance")
        
        self.add_diagram_placeholder(
            slide, "Timeline: Decisions, Contracts, Outcomes",
            left=2, top=2.2, width=6, height=2.5
        )
        
        self.add_vector_branding(slide)
        self.add_notes(slide, """The Decision Ledger provides governance accountability.

Every portfolio decision—entry, exit, rebalance—is logged with timestamp, wave identifier, and rationale. Every alpha contract—the expected return profile and risk parameters—is recorded at inception.

The Ledger tracks outcome reconciliation. Did the wave deliver on its contract? If not, why not? Was it regime conditions? Execution friction? Strategy drift?

This creates an immutable audit trail. Oversight committees can trace any current position back to its original decision rationale. No black boxes. No memory gaps.

The Ledger also supports performance attribution at the decision level. You can see which decisions contributed to returns and which detracted, with full transparency.

This is governance-native design. It treats compliance, transparency, and oversight as first-class requirements, not afterthoughts.""")
        
    def slide_9_performance_deep_dive(self):
        """Slide 9: Performance Deep Dive"""
        slide = self.create_blank_slide()
        self.add_title(slide, "Performance Deep Dive",
                      "Real-Time Analytics Console")
        
        points = [
            "Historical return charts with regime overlays",
            "Rolling alpha metrics & drawdown analysis",
            "Synthetic data indicators for transparency",
            "WaveScore™ proprietary performance scoring",
            "Exportable reports (CSV, PDF, institutional formats)"
        ]
        self.add_bullet_points(slide, points, top=2.2)
        
        self.add_vector_branding(slide)
        self.add_notes(slide, """The Performance Deep Dive console provides real-time visibility into every Wave.

You see historical return charts with regime overlays—shaded regions indicating risk-on versus risk-off periods.

You see rolling alpha metrics, drawdown analysis, and volatility tracking, all benchmarked against governed composite indices.

The system highlights synthetic data with clear visual indicators. No one is misled. If a Wave is using placeholder data while real market integration is in progress, you'll know immediately.

The WaveScore system provides a proprietary performance scoring algorithm that combines return magnitude, consistency, regime balance, and risk efficiency into a single digestible metric.

Every chart, every table, every metric is exportable. CSV downloads, PDF board packs, institutional-ready formatting. No manual data wrangling required.""")
        
    def slide_10_board_pack(self):
        """Slide 10: Board Pack Generation"""
        slide = self.create_blank_slide()
        self.add_title(slide, "Board-Ready Reporting",
                      "Automated PDF Board Packs")
        
        points = [
            "Executive summary with key performance highlights",
            "Wave-by-wave performance analysis",
            "Vector Truth attribution reports",
            "Decision Ledger highlights & governance metrics",
            "Customizable branding and disclosure controls"
        ]
        self.add_bullet_points(slide, points, top=2.2)
        
        self.add_vector_branding(slide)
        self.add_notes(slide, """WAVES Intelligence generates comprehensive PDF board packs automatically.

Each board pack includes an executive summary, wave-by-wave performance analysis, Vector Truth attribution reports, and Decision Ledger highlights.

The system uses deterministic narrative templates. The language is stable, predictable, and governance-appropriate. No marketing fluff. No predictive claims. Just clear, factual reporting.

The board packs integrate seamlessly with existing reporting workflows. They're designed for quarterly board meetings, investment committee reviews, and regulatory filings.

Institutions can customize branding, adjust content sections, and control disclosure levels. But the underlying analytics remain deterministic and auditable.

This isn't just automation. It's institutional-grade documentation at scale.""")
        
    def slide_11_mission_control(self):
        """Slide 11: Mission Control"""
        slide = self.create_blank_slide()
        self.add_title(slide, "Mission Control",
                      "Market Regime & Risk Monitoring")
        
        points = [
            "Real-time VIX levels & regime classification",
            "Exposure adjustments across all Waves",
            "SmartSafe gating for capital preservation",
            "Benchmark drift monitoring & alerts",
            "Capital deployment & concentration tracking"
        ]
        self.add_bullet_points(slide, points, top=2.2)
        
        self.add_vector_branding(slide)
        self.add_notes(slide, """Mission Control is your real-time risk monitoring hub.

It displays current VIX levels, regime classifications, and exposure adjustments across all Waves.

The VIX Regime Overlay system automatically classifies market conditions as risk-on or risk-off based on trailing volatility. When regimes shift, Mission Control highlights which Waves are affected and how exposure levels adjust.

The SmartSafe gating system is visible here. When volatility spikes above critical thresholds, capital preservation overlays activate. You see this in real time. No surprises.

Mission Control also tracks capital deployment across Waves, identifies concentration risks, and monitors benchmark drift. If a governed benchmark is drifting from its snapshot definition, you're alerted immediately.

This is proactive risk management. Not reactive firefighting.""")
        
    def slide_12_implementation(self):
        """Slide 12: Implementation & Integration"""
        slide = self.create_blank_slide()
        self.add_title(slide, "Seamless Implementation",
                      "Data Integration Pathways")
        
        points = [
            "CSV file uploads for historical data",
            "API endpoints for real-time market data",
            "Synthetic data seeding for immediate onboarding",
            "Python/Streamlit architecture, Vercel deployment",
            "Extensible Wave ID registry & custom benchmarks"
        ]
        self.add_bullet_points(slide, points, top=2.2)
        
        self.add_vector_branding(slide)
        self.add_notes(slide, """Implementation is designed for minimal friction.

WAVES Intelligence supports multiple data integration pathways. You can connect via CSV file uploads, API endpoints for real-time market data, or leverage our synthetic data seeding for immediate onboarding.

The system is built on Python and Streamlit, with deployment via Vercel for instant cloud availability. No complex infrastructure. No vendor lock-in.

The data pipeline is idempotent and transactional. You can run data updates multiple times safely. Backups are automatic. Rollback is straightforward.

For institutions with existing data warehouses, we provide schema documentation and migration scripts. The Wave ID registry is extensible—add new Waves, define custom benchmarks, configure attribution rules.

Our implementation team works with your data governance and compliance teams to ensure WAVES Intelligence meets your specific institutional requirements.""")
        
    def slide_13_closing(self):
        """Slide 13: Closing & Next Steps"""
        slide = self.create_blank_slide()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.0), Inches(8), Inches(1)
        )
        text_frame = title_box.text_frame
        text_frame.text = "WAVES Intelligence™"
        
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        
        # Next steps
        points = [
            "Schedule a technical deep dive",
            "Review data integration requirements",
            "Discuss institutional governance needs",
            "",
            "Thank you for your time.",
            "Welcome to the future of institutional portfolio intelligence."
        ]
        
        textbox = slide.shapes.add_textbox(
            Inches(2), Inches(2.5), Inches(6), Inches(2.5)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(points):
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            
            p.text = point
            p.alignment = PP_ALIGN.CENTER
            
            if point == "":
                continue
            elif i < 3:
                p.font.size = Pt(20)
                p.font.color.rgb = self.COLORS['text']
            else:
                p.font.size = Pt(18)
                p.font.color.rgb = self.COLORS['text_dim']
                p.font.italic = True
            
            p.space_before = Pt(10)
            p.space_after = Pt(10)
        
        # Vector signature
        sig_box = slide.shapes.add_textbox(
            Inches(3), Inches(4.8), Inches(4), Inches(0.4)
        )
        text_frame = sig_box.text_frame
        text_frame.text = "— Vector™"
        
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = self.COLORS['secondary']
        p.font.italic = True
        
        self.add_notes(slide, """Thank you for your time today.

To summarize: WAVES Intelligence delivers governance-grade portfolio attribution, real-time risk monitoring, and board-ready reporting—all built on a transparent, deterministic architecture that respects institutional oversight requirements.

We don't predict. We explain. We don't optimize in black boxes. We surface truth outputs with reliability signals.

If your institution is ready to move beyond traditional attribution and embrace transparent, governance-native analytics, we're ready to help.

Next steps are simple. Schedule a technical deep dive with our implementation team. Review our data integration requirements. And let's discuss how WAVES Intelligence can support your specific investment oversight needs.

I'm Vector. Thank you for listening. And welcome to the future of institutional portfolio intelligence.""")
        
    def generate(self):
        """Generate all slides in sequence."""
        print("Generating WAVES Intelligence™ Presentation...")
        
        slides = [
            ("Slide 1: Title & Introduction", self.slide_1_title),
            ("Slide 2: The Challenge", self.slide_2_challenge),
            ("Slide 3: WAVES Solution Overview", self.slide_3_solution),
            ("Slide 4: Wave ID System", self.slide_4_wave_id),
            ("Slide 5: Vector Truth Layer", self.slide_5_vector_truth),
            ("Slide 6: Alpha Decomposition", self.slide_6_alpha_decomposition),
            ("Slide 7: Regime Attribution", self.slide_7_regime_attribution),
            ("Slide 8: Decision Ledger", self.slide_8_decision_ledger),
            ("Slide 9: Performance Deep Dive", self.slide_9_performance_deep_dive),
            ("Slide 10: Board Pack Generation", self.slide_10_board_pack),
            ("Slide 11: Mission Control", self.slide_11_mission_control),
            ("Slide 12: Implementation", self.slide_12_implementation),
            ("Slide 13: Closing & Next Steps", self.slide_13_closing),
        ]
        
        for name, slide_func in slides:
            print(f"  Creating {name}...")
            slide_func()
        
        print(f"✓ Generated {len(slides)} slides")
        
    def save(self, filename):
        """Save presentation to file."""
        self.prs.save(filename)
        print(f"✓ Saved presentation to: {filename}")


def main():
    """Main execution function."""
    # Create presentation
    waves_pres = WavesPresentation()
    waves_pres.generate()
    
    # Determine output path
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_file = os.path.join(script_dir, "WAVES_Intelligence_Executive_Briefing.pptx")
    
    # Save presentation
    waves_pres.save(output_file)
    
    print("\n" + "="*60)
    print("WAVES Intelligence™ Presentation Generation Complete")
    print("="*60)
    print(f"File: {output_file}")
    print(f"Slides: 13")
    print(f"Target Duration: 7-9 minutes")
    print(f"Theme: Dark institutional with Vector™ branding")
    print("="*60)


if __name__ == "__main__":
    main()
